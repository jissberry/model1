#!/usr/bin/env python3
"""生成“极热无风源荷失衡模型”10 页中文科研汇报 PPT。"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "docs" / "figures" / "preview"
EQ = Path("/tmp/ppt-equations")
OUT = ROOT / "docs" / "ppt" / "极热无风源荷失衡模型介绍.pptx"

NAVY = "17365D"
BLUE = "4472C4"
CYAN = "5B9BD5"
RED = "C00000"
ORANGE = "ED7D31"
GOLD = "FFC000"
GREEN = "70AD47"
PURPLE = "7030A0"
DARK = "263238"
GRAY = "667085"
LIGHT = "F4F7FB"
LINE = "D9E2F3"
WHITE = "FFFFFF"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def make_equations():
    """渲染透明公式图片，确保脚本可独立复现。"""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    EQ.mkdir(parents=True, exist_ok=True)
    plt.rcParams["mathtext.fontset"] = "stix"
    items = {
        "source": (
            r"$P_{G,i}^{\max}(T)=P_{G,i}^{\mathrm{rated}}"
            r"[1-\alpha_i(T-T_{\mathrm{ref}})]^{+}$"
        ),
        "load": (
            r"$D_j(T)=\rho_{\mathrm{rigid}}P_{D,j}^{0}"
            r"+\rho_{\mathrm{cool}}P_{D,j}^{0}"
            r"[1+\beta(T-T_{L0})]^{+}$"
        ),
        "opf": (
            r"$\min\ \sum_i(c_{2,i}P_{G,i}^{2}+c_{1,i}P_{G,i})"
            r"+\sum_j\sum_k\mathrm{VOLL}_{k}s_{j,k}$"
        ),
        "pf": r"$P_f=1-\exp(-\lambda\Delta t)$",
        "balance": (
            r"$P_{\mathrm{shed}}^{\star}=(D-P^{\mathrm{avail}})"
            r"+(P^{\mathrm{avail}}-P_G^{\star})$"
        ),
    }
    for name, equation in items.items():
        fig = plt.figure(figsize=(8, 0.65))
        fig.patch.set_alpha(0)
        fig.text(
            0.5, 0.5, equation, ha="center", va="center",
            fontsize=20, color=f"#{NAVY}",
        )
        fig.savefig(
            EQ / f"{name}.png", transparent=True, dpi=220,
            bbox_inches="tight", pad_inches=0.06,
        )
        plt.close(fig)


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def add_text(
    slide, text, x, y, w, h, size=18, color=DARK, bold=False,
    align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font="微软雅黑",
    margin=0.06, line_spacing=1.08,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
    return box


def add_rich_lines(slide, lines, x, y, w, h, size=16, bullet_color=BLUE):
    """lines: [(heading, body), ...]."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    for idx, (heading, body) in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        p.line_spacing = 1.05
        r = p.add_run()
        r.text = "● "
        r.font.name = "微软雅黑"
        r.font.size = Pt(size)
        r.font.color.rgb = rgb(bullet_color)
        r = p.add_run()
        r.text = heading
        r.font.name = "微软雅黑"
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = rgb(NAVY)
        r = p.add_run()
        r.text = body
        r.font.name = "微软雅黑"
        r.font.size = Pt(size)
        r.font.color.rgb = rgb(DARK)
    return box


def add_image_contain(slide, path, x, y, w, h):
    path = Path(path)
    with Image.open(path) as im:
        ratio = im.width / im.height
    box_ratio = w / h
    if ratio > box_ratio:
        pw, ph = w, w / ratio
    else:
        ph, pw = h, h * ratio
    return slide.shapes.add_picture(
        str(path), Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2),
        width=Inches(pw), height=Inches(ph),
    )


def add_header(slide, page, title, section):
    add_rect(slide, 0, 0, 13.333, 0.12, NAVY)
    add_text(slide, f"{page:02d}", 0.45, 0.28, 0.62, 0.42, 20, RED, True)
    add_text(slide, section, 1.1, 0.34, 2.6, 0.3, 10, GRAY, True)
    add_text(slide, title, 0.45, 0.73, 12.2, 0.55, 25, NAVY, True)
    add_rect(slide, 0.45, 1.35, 0.65, 0.045, RED)
    add_rect(slide, 1.1, 1.35, 11.75, 0.018, LINE)
    add_text(slide, f"{page} / 10", 11.9, 7.10, 0.95, 0.2, 8, GRAY, False, PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body, accent=BLUE, title_size=15, body_size=12):
    add_rect(slide, x, y, w, h, WHITE, LINE, True)
    add_rect(slide, x, y, 0.07, h, accent)
    add_text(slide, title, x + 0.18, y + 0.14, w - 0.3, 0.36, title_size, accent, True)
    add_text(slide, body, x + 0.18, y + 0.56, w - 0.34, h - 0.68, body_size, DARK)


def add_kpi(slide, x, y, w, value, label, accent=BLUE):
    add_rect(slide, x, y, w, 0.92, LIGHT, LINE, True)
    value_size = 22
    if len(value) >= 8:
        value_size = 15
    elif len(value) >= 7:
        value_size = 17
    add_text(
        slide, value, x + 0.06, y + 0.10, w - 0.12, 0.38,
        value_size, accent, True, PP_ALIGN.CENTER, margin=0.0,
    )
    add_text(slide, label, x + 0.1, y + 0.53, w - 0.2, 0.22, 9, GRAY, False, PP_ALIGN.CENTER)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(WHITE)
    return slide


def build():
    make_equations()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. 封面：借鉴“科研汇报参考2”的深蓝弧形和“个人展示”的红色强调。
    slide = new_slide(prs)
    add_rect(slide, 0, 0, 13.333, 2.85, NAVY)
    arc = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.8), Inches(1.65), Inches(14.9), Inches(2.25))
    arc.fill.solid()
    arc.fill.fore_color.rgb = rgb(NAVY)
    arc.line.color.rgb = rgb(NAVY)
    add_rect(slide, 0, 3.18, 13.333, 0.035, GOLD)
    add_text(slide, "极热无风条件下", 0.9, 0.68, 11.55, 0.62, 25, WHITE, False, PP_ALIGN.CENTER)
    add_text(slide, "源荷失衡与热故障耦合风险评估", 0.55, 1.28, 12.25, 0.85, 34, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, "基于修改版 IEEE 39 节点系统", 0.9, 2.18, 11.55, 0.36, 15, "D9EAF7", False, PP_ALIGN.CENTER)
    add_text(slide, "模型介绍与算例分析", 4.45, 4.08, 4.45, 0.55, 24, NAVY, True, PP_ALIGN.CENTER)
    add_rect(slide, 5.40, 4.72, 2.55, 0.045, RED)
    add_text(slide, "汇报人：________", 4.65, 5.32, 4.05, 0.35, 14, GRAY, False, PP_ALIGN.CENTER)
    add_text(slide, "2026 年 8 月", 4.65, 5.78, 4.05, 0.35, 13, GRAY, False, PP_ALIGN.CENTER)
    add_text(slide, "SOURCE–LOAD IMBALANCE  ·  THERMAL FAULT  ·  RISK ASSESSMENT",
             1.4, 6.78, 10.53, 0.24, 9, BLUE, True, PP_ALIGN.CENTER, font="Arial")

    # 2. 背景与问题。
    slide = new_slide(prs)
    add_header(slide, 2, "研究背景与问题定义", "01  RESEARCH BACKGROUND")
    add_card(slide, 0.48, 1.65, 3.0, 2.0, "极热天气", "机组冷却效率下降\n光伏温升降效\n设备绝缘老化加速", RED)
    add_card(slide, 3.68, 1.65, 3.0, 2.0, "静风与枯水", "风速低于切入风速\n水电可用容量折减\n新能源支撑能力下降", BLUE)
    add_card(slide, 6.88, 1.65, 3.0, 2.0, "温敏负荷增长", "空调负荷快速攀升\n供需裕度被持续压缩\n形成确定性源荷缺口", PURPLE)
    add_card(slide, 10.08, 1.65, 2.77, 2.0, "热故障叠加", "重载与高温共同抬升\n发电机、线路、变压器\n故障概率", ORANGE)
    add_text(slide, "复合灾害链", 0.55, 4.05, 1.55, 0.36, 16, NAVY, True)
    add_rect(slide, 2.0, 4.19, 10.6, 0.04, LINE)
    chain = [
        ("气象冲击", NAVY), ("源荷缺口", BLUE), ("潮流重分布", GREEN),
        ("热故障概率升高", ORANGE), ("故障后切负荷", RED),
    ]
    for i, (label, color) in enumerate(chain):
        x = 0.85 + i * 2.48
        add_rect(slide, x, 4.65, 1.82, 0.72, color, color, True)
        add_text(slide, label, x + 0.05, 4.84, 1.72, 0.28, 12, WHITE, True, PP_ALIGN.CENTER)
        if i < len(chain) - 1:
            add_text(slide, "→", x + 1.88, 4.79, 0.55, 0.34, 22, GRAY, True, PP_ALIGN.CENTER)
    add_rect(slide, 0.7, 5.80, 11.95, 0.95, LIGHT, LINE, True)
    add_text(slide, "核心问题", 0.92, 6.02, 1.25, 0.3, 14, RED, True)
    add_text(slide, "如何在统一框架中量化“确定性源荷失衡”与“随机元件热故障”对失电风险的耦合影响？",
             2.08, 5.96, 10.15, 0.42, 16, NAVY, True)
    add_text(slide, "背景数据来源：文献[1]–[17]；本汇报数值均为修改版 IEEE 39 节点算例结果。",
             0.7, 6.84, 10.7, 0.2, 8, GRAY)

    # 3. 总体框架。
    slide = new_slide(prs)
    add_header(slide, 3, "研究思路与总体框架", "02  METHODOLOGY")
    add_image_contain(slide, FIG / "preview_01_framework_v2.png", 0.48, 1.52, 12.36, 5.30)
    add_text(slide, "两阶段运行逻辑：基准 OPF 无爬坡约束；故障后 OPF 以基准出力为爬坡中心。",
             1.4, 6.80, 10.55, 0.24, 10, RED, True, PP_ALIGN.CENTER)

    # 4. 源荷失衡模型。
    slide = new_slide(prs)
    add_header(slide, 4, "极热无风条件下的源荷失衡模型", "02  METHODOLOGY")
    add_image_contain(slide, FIG / "preview_02_source_load_v2.png", 0.45, 1.53, 12.4, 3.70)
    add_image_contain(slide, EQ / "source.png", 0.62, 5.30, 5.85, 0.58)
    add_image_contain(slide, EQ / "load.png", 6.72, 5.30, 5.95, 0.58)
    add_text(slide, "源侧", 0.7, 5.91, 0.65, 0.26, 12, RED, True)
    add_text(slide, "火电高温降容 · 水电枯水折减 · 风电低风速停发 · 光伏温度—辐照修正",
             1.35, 5.89, 5.15, 0.43, 10, DARK)
    add_text(slide, "荷侧", 6.8, 5.91, 0.65, 0.26, 12, PURPLE, True)
    add_text(slide, "刚性负荷 + 温敏负荷；按一级/二级/三级设置差异化 VOLL",
             7.45, 5.89, 5.15, 0.43, 10, DARK)
    add_rect(slide, 0.68, 6.49, 11.95, 0.42, "FFF2F2", "F4CCCC", True)
    add_text(slide, "40 °C、2 m/s 场景：可用出力下降与负荷增长共同形成 746.21 MW 原始能量缺口。",
             0.85, 6.56, 11.55, 0.23, 12, RED, True, PP_ALIGN.CENTER)

    # 5. 优化模型和算例。
    slide = new_slide(prs)
    add_header(slide, 5, "基准 DC-OPF / MIQP 与算例设置", "03  OPTIMIZATION MODEL")
    add_card(slide, 0.48, 1.58, 7.65, 1.25, "优化目标：发电经济性 + 分级失负荷代价", "", BLUE)
    add_image_contain(slide, EQ / "opf.png", 0.75, 2.02, 7.05, 0.62)
    add_card(slide, 0.48, 3.02, 3.67, 1.48, "网络约束", "节点功率平衡\n直流支路潮流\n支路热稳限额与平衡节点", GREEN)
    add_card(slide, 4.46, 3.02, 3.67, 1.48, "机组约束", "火电启停 u∈{0,1}\n极热修正出力上下界\n基准状态不施加爬坡", ORANGE)
    add_card(slide, 0.48, 4.72, 7.65, 1.48, "分级切负荷", "VOLL₁=10000、VOLL₂=5000、VOLL₃=1000 $/MWh\n优化模型优先保障一级关键负荷和二级重要负荷。", PURPLE)
    add_rect(slide, 8.48, 1.58, 4.37, 4.62, LIGHT, LINE, True)
    add_text(slide, "修改版 IEEE 39 节点系统", 8.78, 1.83, 3.75, 0.38, 18, NAVY, True, PP_ALIGN.CENTER)
    add_kpi(slide, 8.83, 2.45, 1.62, "39", "节点", NAVY)
    add_kpi(slide, 10.58, 2.45, 1.62, "46", "支路", BLUE)
    add_kpi(slide, 8.83, 3.55, 1.62, "10", "发电机", ORANGE)
    add_kpi(slide, 10.58, 3.55, 1.62, "29", "节点变压器", PURPLE)
    add_text(slide, "气象输入", 8.78, 4.78, 1.15, 0.3, 13, RED, True)
    add_text(slide, "T = 40 °C\nv = 2 m/s\nG = 900 W/m²",
             9.78, 4.70, 2.25, 1.05, 15, DARK, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_text(slide, "实现：MATLAB/Gurobi 主模型 + Python 连续 QP 枚举交叉验证",
             8.72, 5.76, 3.85, 0.34, 9, GRAY, False, PP_ALIGN.CENTER)

    # 6. 基准结果。
    slide = new_slide(prs)
    add_header(slide, 6, "基准 OPF：能量缺口与网络约束", "04  BASELINE RESULTS")
    add_image_contain(slide, FIG / "preview_03_dispatch_v2.png", 0.4, 1.55, 8.65, 4.82)
    add_kpi(slide, 9.33, 1.72, 1.45, "6754.57", "极热需求 / MW", PURPLE)
    add_kpi(slide, 10.92, 1.72, 1.45, "6008.36", "可用上界 / MW", BLUE)
    add_kpi(slide, 9.33, 2.82, 1.45, "746.21", "原始缺口 / MW", ORANGE)
    add_kpi(slide, 10.92, 2.82, 1.45, "762.81", "最优切负荷 / MW", RED)
    add_card(slide, 9.20, 4.05, 3.40, 1.45, "为什么多出 16.60 MW？",
             "G4 极热上界 668.5 MW\n受支路热稳约束仅出力 651.9 MW\n差值即网络约束闲置能力。", RED, 14, 11)
    add_image_contain(slide, EQ / "balance.png", 9.12, 5.62, 3.62, 0.55)
    add_rect(slide, 0.65, 6.46, 11.95, 0.42, "EFF8EF", "C6E0B4", True)
    add_text(slide, "结果：总切负荷占需求 11.29%，全部来自三级可中断负荷；线路最大负载率 100%，无越限。",
             0.86, 6.54, 11.55, 0.22, 11, GREEN, True, PP_ALIGN.CENTER)

    # 7. 故障概率。
    slide = new_slide(prs)
    add_header(slide, 7, "极热条件下的元件热故障概率", "05  THERMAL FAULT MODEL")
    add_image_contain(slide, FIG / "preview_05_fault_probability_v2.png", 0.38, 1.56, 9.45, 4.58)
    add_card(slide, 10.03, 1.62, 2.80, 1.15, "统一概率映射", "", NAVY, 14, 11)
    add_image_contain(slide, EQ / "pf.png", 10.20, 2.02, 2.46, 0.51)
    add_card(slide, 10.03, 2.98, 2.80, 1.20, "节点变压器", "热点温度 + Arrhenius\n绝缘热老化加速", PURPLE, 13, 10)
    add_card(slide, 10.03, 4.32, 2.80, 1.20, "发电机", "温度应力 + 出力率\n比例风险模型", ORANGE, 13, 10)
    add_card(slide, 10.03, 5.66, 2.80, 1.20, "输电线路", "IEEE 738 导线温度\n温度/过载应力放大", RED, 13, 10)
    add_text(slide, "24 h 条件故障概率估计：G10≈4.11%，G4≈4.04%；最高风险线路 L27≈0.42%；变压器≈0.015%。",
             0.65, 6.46, 9.05, 0.44, 10, NAVY, True, PP_ALIGN.CENTER)

    # 8. MC 与故障后 OPF。
    slide = new_slide(prs)
    add_header(slide, 8, "蒙特卡洛场景与故障后 OPF", "06  SCENARIO ANALYSIS")
    add_image_contain(slide, FIG / "preview_06_monte_carlo_v2.png", 0.42, 1.62, 7.22, 3.00)
    add_card(slide, 7.95, 1.62, 4.88, 1.02, "① 85 维故障向量",
             "10 台发电机 + 46 条线路 + 29 台节点变压器", NAVY, 14, 10)
    add_card(slide, 7.95, 2.80, 4.88, 1.02, "② 0/1 故障映射",
             "机组退出；线路退出；变压器故障映射为关联支路退出", PURPLE, 14, 10)
    add_card(slide, 7.95, 3.98, 4.88, 1.02, "③ 故障后重调度",
             "以基准 Pᴳ★ 为爬坡中心；火电可停机处理孤岛最小出力过剩", RED, 14, 10)
    add_kpi(slide, 0.62, 5.18, 1.78, "2000", "场景数", BLUE)
    add_kpi(slide, 2.58, 5.18, 1.78, "739", "总故障事件", ORANGE)
    add_kpi(slide, 4.54, 5.18, 1.78, "631", "含故障场景", PURPLE)
    add_kpi(slide, 6.50, 5.18, 1.78, "70.2%", "发电机事件占比", RED)
    add_kpi(slide, 8.46, 5.18, 1.78, "388", "含火电关停场景", GREEN)
    add_kpi(slide, 10.42, 5.18, 1.78, "2000/2000", "获得最优解", NAVY)
    add_rect(slide, 0.68, 6.38, 11.80, 0.48, LIGHT, LINE, True)
    add_text(slide, "抽样一致性：经验故障频率与理论概率总体贴合；单场景最多 3 个故障元件。",
             0.85, 6.49, 11.45, 0.22, 11, NAVY, True, PP_ALIGN.CENTER)

    # 9. 风险分布与驱动因素。
    slide = new_slide(prs)
    add_header(slide, 9, "故障后切负荷分布与风险主导因素", "07  RISK RESULTS")
    add_image_contain(slide, FIG / "preview_07_shed_distribution_v2.png", 0.35, 1.55, 7.45, 2.75)
    add_image_contain(slide, FIG / "preview_08_risk_drivers_v2.png", 0.35, 4.36, 7.45, 2.55)
    add_kpi(slide, 8.12, 1.64, 1.42, "951.58", "均值 / MW", BLUE)
    add_kpi(slide, 9.68, 1.64, 1.42, "1765.56", "P95 / MW", ORANGE)
    add_kpi(slide, 11.24, 1.64, 1.42, "3330.40", "最大值 / MW", RED)
    add_card(slide, 8.08, 2.82, 4.65, 1.24, "分布特征",
             "约 75% 场景仍停留在 762.81 MW 基准切负荷；风险主要集中于右侧尾部。", BLUE, 14, 11)
    add_card(slide, 8.08, 4.24, 4.65, 1.24, "主导因素",
             "电源故障后果强于单纯网络故障；电源与网络混合故障对应最高风险。", RED, 14, 11)
    add_card(slide, 8.08, 5.66, 4.65, 1.24, "尾部口径",
             "90% 分位总切负荷 1542.81 MW，相对基准新增 780 MW；最大值约为基准 4.4 倍。", ORANGE, 14, 11)

    # 10. 结论、边界和展望。
    slide = new_slide(prs)
    add_header(slide, 10, "主要结论、模型边界与展望", "08  CONCLUSION")
    add_text(slide, "主要结论", 0.52, 1.58, 2.0, 0.40, 18, NAVY, True)
    add_rich_lines(slide, [
        ("确定性缺口：", "40 °C 极热无风形成约 746 MW 原始能量缺口；考虑网络约束后切负荷约 763 MW。"),
        ("关键负荷保障：", "差异化 VOLL 将基准切负荷全部集中在三级可中断负荷。"),
        ("尾部风险突出：", "多数场景不增加切负荷，但 P95 达 1766 MW；电源故障是主要风险驱动。"),
        ("方法可复现：", "基准 OPF—故障概率—MC 抽样—故障后 OPF 的数值链条闭合，并经 Python 交叉验证。"),
    ], 0.55, 2.05, 6.05, 3.95, 15, BLUE)
    add_rect(slide, 6.85, 1.58, 5.98, 2.52, LIGHT, LINE, True)
    add_text(slide, "模型边界", 7.12, 1.82, 1.6, 0.35, 17, RED, True)
    add_text(slide,
             "• DC 潮流未计及电压、无功与变压器视在功率\n"
             "• 单时段静态快照；基准无爬坡、故障后有爬坡\n"
             "• 元件独立抽样，未描述区域共因故障与保护动态\n"
             "• 故障概率参数未经真实运维数据标定\n"
             "• 修改版测试系统结果不能直接外推至真实电网",
             7.10, 2.25, 5.38, 1.63, 12, DARK)
    add_rect(slide, 6.85, 4.36, 5.98, 1.95, "EFF5FB", "BDD7EE", True)
    add_text(slide, "下一步工作", 7.12, 4.60, 1.6, 0.35, 17, BLUE, True)
    add_text(slide,
             "① 引入 AC-OPF，刻画电压/无功与视在功率\n"
             "② 扩展多时段调度、储能和需求响应\n"
             "③ 结合真实气象场与设备运维数据标定参数\n"
             "④ 建模共因故障、保护动作与恢复过程",
             7.10, 5.02, 5.38, 1.12, 12, DARK)
    add_rect(slide, 0.65, 6.48, 12.0, 0.44, NAVY, NAVY, True)
    add_text(slide, "所提框架用于极端气象风险评估，而非对特定真实电网的确定性运行预测。",
             0.85, 6.57, 11.6, 0.23, 12, WHITE, True, PP_ALIGN.CENTER)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"written {OUT}")


if __name__ == "__main__":
    build()
