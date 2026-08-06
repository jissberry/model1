#!/usr/bin/env python3
"""科研经历 2.1 研究背景（单页）。

版式对齐《个人展示ppt模板.pptx》第 4 页“2.1 研究背景”：
顶部蓝框背景卡 + 中部内容相关配图 + 底部蓝框意义卡（红字强调）。
配图为自绘的“极热天气下源荷失衡与热故障耦合机理”示意图，与正文直接相关。
"""

from __future__ import annotations

from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FIGPATH = ROOT / "docs" / "figures" / "研究背景_源荷失衡与热故障耦合机理.png"
OUT = ROOT / "docs" / "ppt" / "科研经历_2.1_研究背景.pptx"

BLUE = "0070C0"
DEEP = "1F4E79"
RED = "C00000"
PURPLE = "7030A0"
ORANGE = "C55A11"
TEAL = "2E8B8B"
BLACK = "262626"
GRAY = "595959"
BOXLINE = "2E75B6"
WHITE = "FFFFFF"

TOPIC = "考虑热故障及源荷失衡的极热天气电网失电风险评估"

# matplotlib 颜色
M_HEAT = "#C00000"
M_BLUE = "#1F6FB2"
M_PURPLE = "#7030A0"
M_ORANGE = "#C55A11"
M_TEAL = "#2E8B8B"
M_GRAY = "#5A6472"
M_RISK = "#B01515"


def make_figure() -> None:
    plt.rcParams["font.sans-serif"] = ["WenQuanYi Micro Hei", "Droid Sans Fallback"]
    plt.rcParams["axes.unicode_minus"] = False

    fig, ax = plt.subplots(figsize=(12.2, 3.65))
    ax.set_xlim(0, 122)
    ax.set_ylim(0, 36)
    ax.axis("off")

    def box(x0, x1, y0, y1, color, lw=1.8, fill="#FFFFFF", alpha=1.0, rounding=0.10):
        ax.add_patch(FancyBboxPatch(
            (x0, y0), x1 - x0, y1 - y0,
            boxstyle=f"round,pad=0.02,rounding_size={rounding}",
            linewidth=lw, edgecolor=color, facecolor=fill, alpha=alpha, zorder=2,
        ))

    def arrow(p0, p1, color=M_GRAY, lw=2.0, style="-|>", rad=0.0):
        ax.add_patch(FancyArrowPatch(
            p0, p1, arrowstyle=style, mutation_scale=16,
            linewidth=lw, color=color, zorder=3,
            connectionstyle=f"arc3,rad={rad}",
        ))

    # 顶部：气象冲击
    box(36, 86, 30.2, 35.4, M_HEAT, 2.0, "#FBE9E7")
    ax.text(61, 32.8, "极端高温 · 静风 · 强辐照 · 枯水",
            ha="center", va="center", fontsize=13.5, color=M_HEAT, fontweight="bold")

    # 左：源荷失衡
    box(3.5, 53, 10.5, 27.5, M_BLUE, 2.0, "#EAF3FA")
    ax.text(28.25, 25.6, "源荷失衡", ha="center", va="center",
            fontsize=13.5, color=M_BLUE, fontweight="bold")
    ax.text(28.25, 21.6,
            "源侧降容：火电/水电降容、风电停发、光伏降效  →  可用出力 ↓",
            ha="center", va="center", fontsize=9.6, color="#20364A")
    ax.text(28.25, 18.2,
            "荷侧激增：空调等温敏负荷快速攀升  →  用电需求 ↑",
            ha="center", va="center", fontsize=9.6, color="#20364A")
    ax.add_patch(FancyBboxPatch((11.5, 11.2), 33.5, 2.9,
                 boxstyle="round,pad=0.02,rounding_size=0.12",
                 linewidth=1.4, edgecolor=M_BLUE, facecolor="#D6E7F5", zorder=2))
    ax.text(28.25, 12.6, "⇒  确定性供需缺口",
            ha="center", va="center", fontsize=10.2, color=M_BLUE, fontweight="bold")

    # 右：元件热故障
    box(69, 118.5, 10.5, 27.5, M_ORANGE, 2.0, "#FDF0E6")
    ax.text(93.75, 25.6, "元件热故障", ha="center", va="center",
            fontsize=13.5, color=M_ORANGE, fontweight="bold")
    for i, t in enumerate([
        "变压器：热点温度升高、绝缘热老化加速",
        "输电线路：导线温度升高、过载风险增大",
        "发电机：高温叠加高出力率，停运率上升",
    ]):
        ax.text(93.75, 22.4 - i * 2.9, t, ha="center", va="center",
                fontsize=9.6, color="#4A2E12")
    ax.add_patch(FancyBboxPatch((77, 11.2), 33.5, 2.9,
                 boxstyle="round,pad=0.02,rounding_size=0.12",
                 linewidth=1.4, edgecolor=M_ORANGE, facecolor="#F6DBC4", zorder=2))
    ax.text(93.75, 12.6, "⇒  故障概率显著升高", ha="center", va="center",
            fontsize=10.2, color=M_ORANGE, fontweight="bold")

    # 顶部向两支的箭头
    arrow((50, 30.2), (34, 27.6), M_GRAY, 2.0, rad=0.12)
    arrow((72, 30.2), (88, 27.6), M_GRAY, 2.0, rad=-0.12)

    # 中部耦合双箭头
    arrow((53.4, 19.0), (68.6, 19.0), M_RISK, 2.4, style="<|-|>")
    ax.text(61, 22.4, "耦合放大", ha="center", va="center", fontsize=10.5,
            color=M_RISK, fontweight="bold")
    ax.text(61, 16.0, "①重载→故障↑", ha="center", va="center",
            fontsize=8.2, color=M_RISK)
    ax.text(61, 13.6, "②故障→失衡↑", ha="center", va="center",
            fontsize=8.2, color=M_RISK)

    # 底部风险
    box(33, 89, 1.5, 7.6, M_RISK, 2.0, "#FBE3E3")
    ax.text(61, 4.5, "电网失电风险：大面积停电 · 大规模缺电",
            ha="center", va="center", fontsize=13, color=M_RISK, fontweight="bold")

    # 两支向风险的箭头
    arrow((28.25, 9.5), (48, 7.7), M_BLUE, 2.0, rad=0.10)
    arrow((93.75, 9.5), (74, 7.7), M_ORANGE, 2.0, rad=-0.10)

    fig.savefig(FIGPATH, dpi=240, bbox_inches="tight", pad_inches=0.06, facecolor="white")
    plt.close(fig)


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def _style_run(run, size, color, bold, font):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    rpr = run.font._rPr
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rpr.find(qn(tag))
        if el is None:
            el = rpr.makeelement(qn(tag), {})
            rpr.append(el)
        el.set("typeface", font)


def add_rich(slide, paragraphs, x, y, w, h, size=14, base=BLACK,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, line_spacing=1.25,
             font="微软雅黑", first_indent=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.10)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = valign
    for pi, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for seg in para:
            txt, emph = seg if isinstance(seg, tuple) else (seg, False)
            r = p.add_run()
            r.text = txt
            _style_run(r, size, RED if emph else base, bool(emph), font)
    return box


def add_text(slide, text, x, y, w, h, size, color, bold=False,
             align=PP_ALIGN.LEFT, font="微软雅黑"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _style_run(r, size, color, bold, font)
    return box


def add_line(slide, x, y, w, color=BLUE, weight=2.0):
    ln = slide.shapes.add_connector(2, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = rgb(color)
    ln.line.width = Pt(weight)
    return ln


def add_round(slide, x, y, w, h, line=BOXLINE, weight=1.4, fill=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(weight)
    shape.adjustments[0] = 0.10
    shape.shadow.inherit = False
    return shape


def add_image_contain(slide, path, x, y, w, h):
    with Image.open(path) as im:
        ratio = im.width / im.height
    if ratio > w / h:
        pw, ph = w, w / ratio
    else:
        ph, pw = h, h * ratio
    return slide.shapes.add_picture(str(path),
                                    Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2),
                                    width=Inches(pw), height=Inches(ph))


def build_slide() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(WHITE)

    # 主标题 + 分隔线
    add_text(slide, "二、科研经历—" + TOPIC, 0.35, 0.12, 12.8, 0.5, 22, BLUE, True)
    add_line(slide, 0.0, 0.74, 13.333, BLUE, 2.2)
    add_text(slide, "2.1  研究背景", 0.35, 0.80, 6.0, 0.42, 18, BLUE, True)

    # 顶部背景卡
    add_round(slide, 0.55, 1.34, 12.23, 1.16)
    add_rich(slide, [
        [("近年来，全球", False), ("极端高温、静风", True),
         ("等复合气象事件频发。2022 年夏季四川持续极端高温，日最大电力缺口约 ", False),
         ("1700 万千瓦", True),
         ("；2021 年美国得州极端天气引发大停电，逾 ", False),
         ("330 万用户", True), ("失去正常供电。", False)],
        [("极热天气通过", False), ("源侧降容、荷侧激增与元件热故障", True),
         ("等多条路径，持续压缩电网安全裕度。", False)],
    ], 0.75, 1.42, 11.85, 1.0, size=13.5, align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE, line_spacing=1.3)

    # 中部配图
    add_image_contain(slide, FIGPATH, 0.55, 2.72, 12.23, 3.42)
    add_text(slide, "图 2-1  极热天气下源荷失衡与热故障耦合致失电风险机理",
             0.55, 6.12, 12.23, 0.28, 10.5, GRAY, False, PP_ALIGN.CENTER)

    # 底部意义卡
    add_round(slide, 1.02, 6.44, 11.30, 0.92)
    add_rich(slide, [
        [("极热天气下", False), ("源荷失衡", True), ("与", False),
         ("元件热故障", True), ("相互耦合、彼此放大；准确评估二者耦合下的", False),
         ("电网失电风险", True),
         ("，是极端气象下保障可靠供电与制定韧性提升策略的前提。", False)],
    ], 1.22, 6.50, 10.9, 0.80, size=13.5, align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE, line_spacing=1.25)

    add_text(slide, "1 / 6", 12.35, 7.08, 0.85, 0.24, 9, GRAY, False, PP_ALIGN.RIGHT)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"written {OUT}")


if __name__ == "__main__":
    make_figure()
    build_slide()
