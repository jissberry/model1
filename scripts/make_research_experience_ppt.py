#!/usr/bin/env python3
"""生成“科研经历”模块 PPT（6 页）。

主题：考虑热故障及源荷失衡的极热天气电网失电风险评估。
版式严格对齐仓库根目录《个人展示ppt模板.pptx》的“二、科研经历”系列：
蓝色主标题 + 分隔线 + “2.X 小标题” + 蓝框白底文本卡（红字强调）+ 居中公式 + 图 + 蓝色箭头。
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "docs" / "figures" / "preview"
EQ = Path("/tmp/research-ppt-eq")
OUT = ROOT / "docs" / "ppt" / "科研经历_极热天气电网失电风险评估.pptx"

BLUE = "0070C0"      # 模板主蓝
DEEP = "1F4E79"      # 深蓝
RED = "C00000"       # 强调红
BLACK = "262626"
GRAY = "595959"
LIGHTBLUE = "DEEBF7"
BOXLINE = "2E75B6"
WHITE = "FFFFFF"

TOPIC = "考虑热故障及源荷失衡的极热天气电网失电风险评估"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def make_equations() -> None:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    EQ.mkdir(parents=True, exist_ok=True)
    plt.rcParams["mathtext.fontset"] = "stix"
    items = {
        "derate": r"$P_{G,i}^{\max}(T)=P_{G,i}^{\mathrm{rated}}\left[1-\alpha_i\left(T-T_{\mathrm{ref}}\right)\right]^{+}$",
        "load": r"$D_j(T)=\rho_{\mathrm{r}}P_{D,j}^{0}+\rho_{\mathrm{c}}P_{D,j}^{0}\left[1+\beta\left(T-T_{L0}\right)\right]^{+}$",
        "obj": r"$\min\ \sum_{i}\left(c_{2,i}P_{G,i}^{2}+c_{1,i}P_{G,i}\right)+\sum_{j}\sum_{k=1}^{3}\mathrm{VOLL}_{k}\,s_{j,k}$",
        "balance": r"$S_B(\mathbf{B}\boldsymbol{\theta})_n-\sum_{i\in\mathcal{G}(n)}P_{G,i}-\sum_{k}s_{n,k}=-D_n(T)$",
        "ucon": r"$u_i \underline{P}_{G,i} \leq P_{G,i} \leq u_i P_{G,i}^{\max}(T),\ \ u_i \in \{0,1\}$",
        "pf": r"$P_f=1-\exp\!\left(-\lambda\,\Delta t\right)$",
        "trafo": r"$\lambda_T=\lambda_{T,0}\exp\!\left(\dfrac{B}{\theta_{H,\mathrm{ref}}+273}-\dfrac{B}{\theta_H+273}\right)$",
        "gen": r"$\lambda_{G,i}=\lambda_{0}\exp\!\left[a_T\!\left(T_{\mathrm{eff},i}-T_0\right)+a_L\,\ell_i\right]$",
        "line": r"$\lambda_{L,l}=\lambda_{L,0}\exp\!\left[b_T(T_{c,l}-T_{c,\mathrm{ref}})^{+}+b_S(\beta_l-1)^{+}\right]$",
        "sample": r"$z_i^{(s)}=0\ (u_i^{(s)}<p_i)\ ;\quad z_i^{(s)}=1\ (u_i^{(s)}\geq p_i)$",
        "ramp": r"$\overline{P}_{G,i,s}^{\mathrm{disp}}=\min\!\left(P_{G,i}^{\max}(T),\,P_{G,i}^{\star}+R_i^{\uparrow}\Delta t\right)$",
        "bridge": r"$P_{\mathrm{shed}}^{\star}=(D-P^{\mathrm{avail}})+(P^{\mathrm{avail}}-P_G^{\star})$",
    }
    for name, equation in items.items():
        fig = plt.figure(figsize=(8, 0.7))
        fig.patch.set_alpha(0)
        fig.text(0.5, 0.5, equation, ha="center", va="center", fontsize=19, color=f"#{DEEP}")
        fig.savefig(EQ / f"{name}.png", transparent=True, dpi=230, bbox_inches="tight", pad_inches=0.05)
        plt.close(fig)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(WHITE)
    return slide


def add_line(slide, x, y, w, color=BLUE, weight=2.0):
    ln = slide.shapes.add_connector(2, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = rgb(color)
    ln.line.width = Pt(weight)
    return ln


def add_rect(slide, x, y, w, h, fill=None, line=None, radius=False, weight=1.25):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(weight)
    if radius:
        shape.adjustments[0] = 0.06
    shape.shadow.inherit = False
    return shape


def _style_run(run, size, color, bold, font):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    ea = run.font._rPr.get_or_change_to_latin() if False else None  # noqa
    rpr = run.font._rPr
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rpr.find(qn(tag))
        if el is None:
            el = rpr.makeelement(qn(tag), {})
            rpr.append(el)
        el.set("typeface", font)


def add_text(slide, text, x, y, w, h, size=14, color=BLACK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font="微软雅黑",
             line_spacing=1.12, margin=0.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        _style_run(run, size, color, bold, font)
    return box


def add_rich(slide, segments, x, y, w, h, size=14, base=BLACK,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, line_spacing=1.2,
             bullet=False, bullet_color=RED, font="微软雅黑"):
    """segments: list of paragraphs; each paragraph is list of (text, emphasis?)."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = valign
    for pi, para in enumerate(segments):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(5)
        if bullet:
            r = p.add_run()
            r.text = "■ "
            _style_run(r, size, bullet_color, True, font)
        for seg in para:
            txt, emph = seg if isinstance(seg, tuple) else (seg, False)
            r = p.add_run()
            r.text = txt
            _style_run(r, size, RED if emph else base, bool(emph), font)
    return box


def add_image_contain(slide, path, x, y, w, h):
    with Image.open(path) as im:
        ratio = im.width / im.height
    if ratio > w / h:
        pw, ph = w, w / ratio
    else:
        ph, pw = h, h * ratio
    return slide.shapes.add_picture(
        str(path), Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2),
        width=Inches(pw), height=Inches(ph),
    )


def add_arrow(slide, x, y, w, h, color=BLUE):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid()
    a.fill.fore_color.rgb = rgb(color)
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def add_dash_divider(slide, x, y, h, color=BOXLINE):
    ln = slide.shapes.add_connector(2, Inches(x), Inches(y), Inches(x), Inches(y + h))
    ln.line.color.rgb = rgb(color)
    ln.line.width = Pt(1.1)
    ln.line._get_or_add_ln().append(
        ln.line._get_or_add_ln().makeelement(qn("a:prstDash"), {"val": "dash"})
    )
    return ln


def header(slide, sub_no, sub_title, page):
    add_rich(
        slide,
        [[("二、科研经历", False), ("—" + TOPIC, False)]],
        0.35, 0.12, 12.7, 0.5, size=22, base=BLUE, align=PP_ALIGN.LEFT,
    )
    for shp in slide.shapes:
        pass
    # 主标题加粗蓝
    add_line(slide, 0.0, 0.74, 13.333, BLUE, 2.2)
    add_rich(
        slide,
        [[(sub_no + "  ", True), (sub_title, False)]],
        0.35, 0.80, 12.6, 0.42, size=17, base=BLUE, align=PP_ALIGN.LEFT,
    )
    add_text(slide, f"{page} / 6", 12.35, 7.06, 0.85, 0.24, 9, GRAY, False, PP_ALIGN.RIGHT)


def bold_title(slide):
    # 让主标题真正加粗（add_rich 已 bold=False，这里覆盖首个 textbox 的 runs）
    tb = slide.shapes[0]
    for p in tb.text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True


def callout(slide, x, y, w, h, paragraphs, size=13, valign=MSO_ANCHOR.MIDDLE, bullet=False):
    add_rect(slide, x, y, w, h, WHITE, BOXLINE, True, 1.4)
    add_rich(slide, paragraphs, x + 0.12, y + 0.05, w - 0.24, h - 0.1, size=size,
             valign=valign, bullet=bullet, line_spacing=1.18)


def label(slide, text, x, y, w=3.4, size=13, color=BLUE):
    add_rich(slide, [[("· ", True), (text, False)]], x, y, w, 0.32, size=size, base=color)


def build():
    make_equations()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---------- 2.1 研究背景与总体思路 ----------
    s = new_slide(prs)
    header(s, "2.1", "研究背景与总体思路", 1)
    bold_title(s)
    callout(s, 0.55, 1.30, 12.25, 0.98, [[
        ("近年来极端高温、静风等复合气象事件频发：", False),
        ("源侧机组降容、荷侧空调负荷激增", True),
        ("，同时高温重载抬升发电机、线路与变压器的", False),
        ("热故障概率", True),
        ("。亟需在统一框架下评估源荷失衡与热故障耦合下的电网失电风险。", False),
    ]], size=13.5)
    add_image_contain(s, FIG / "preview_01_framework_v2.png", 0.5, 2.45, 8.35, 4.55)
    label(s, "研究流程（两阶段运行逻辑）", 9.1, 2.45, 4.0, 13, DEEP)
    steps = [
        ("① 源荷失衡建模", "四类电源降容 + 温敏负荷增长", BLUE),
        ("② 基准 DC-OPF", "MIQP 求解；不施加爬坡约束", DEEP),
        ("③ 热故障概率", "变压器 / 发电机 / 线路", RED),
        ("④ 蒙特卡洛场景", "2000 × 85 维故障抽样", BLUE),
        ("⑤ 故障后 OPF", "以基准出力为爬坡中心", DEEP),
        ("⑥ 失电风险统计", "均值 / 分位数 / 尾部风险", RED),
    ]
    y = 2.86
    for t, b, c in steps:
        add_rect(s, 9.1, y, 3.75, 0.60, LIGHTBLUE, BOXLINE, True, 1.0)
        add_rich(s, [[(t, False)]], 9.22, y + 0.03, 3.6, 0.28, size=12.5, base=c)
        add_text(s, b, 9.22, y + 0.31, 3.55, 0.24, 9.5, GRAY)
        y += 0.685
    add_text(s, "算例：修改版 IEEE 39 节点系统（39 节点 / 46 支路 / 10 机 / 29 节点变压器）",
             0.55, 7.05, 8.3, 0.26, 10, GRAY)

    # ---------- 2.2 极热无风源荷失衡建模 ----------
    s = new_slide(prs)
    header(s, "2.2", "极热无风条件下的源荷失衡建模", 2)
    bold_title(s)
    add_dash_divider(s, 6.72, 1.35, 5.35)
    label(s, "源侧：四类电源可用出力修正", 0.4, 1.30, 6.2, 13, DEEP)
    add_image_contain(s, EQ / "derate.png", 0.45, 1.66, 6.0, 0.52)
    add_rich(s, [[
        ("火电按", False), ("进气密度、冷却效率", True),
        ("下降线性降容；水电受", False), ("枯水折减", True),
        ("；风速低于切入风速时", False), ("风电近乎停发", True),
        ("；光伏按", False), ("温度—辐照", True), ("修正。", False),
    ]], 0.42, 2.22, 6.2, 0.9, size=12.5)
    add_image_contain(s, FIG / "preview_02_source_load_v2.png", 0.30, 3.05, 6.45, 3.55)
    label(s, "荷侧：温敏负荷增长与分级切负荷", 6.95, 1.30, 6.2, 13, DEEP)
    add_image_contain(s, EQ / "load.png", 6.9, 1.66, 6.1, 0.52)
    add_rich(s, [[
        ("负荷分为", False), ("刚性 + 温敏", True),
        ("两部分，温敏负荷随温度上升；并按", False),
        ("一级 / 二级 / 三级", True),
        ("设置差异化失负荷价值 VOLL，优先保障关键负荷。", False),
    ]], 6.92, 2.22, 6.1, 0.9, size=12.5)
    callout(s, 6.95, 3.20, 5.9, 1.55, [
        [("源荷缺口的形成：", True)],
        [("• 40 °C、2 m/s 场景下，系统可用出力上界约 ", False), ("6008.36 MW", True)],
        [("• 极热修正后总需求约 ", False), ("6754.57 MW", True)],
        [("• 二者形成确定性", False), ("原始能量缺口 ≈ 746.21 MW", True)],
    ], size=12.5, valign=MSO_ANCHOR.MIDDLE)
    callout(s, 6.95, 4.95, 5.9, 1.65, [
        [("分级切负荷约束：", True)],
        [("0 ≤ s(j,k) ≤ D(j,k)(T)，  k = 1,2,3", False)],
        [("VOLL₁ = 10000，VOLL₂ = 5000，VOLL₃ = 1000 $/MWh", False)],
        [("→ 引导模型把切负荷集中到", False), ("三级可中断负荷", True)],
    ], size=12, valign=MSO_ANCHOR.MIDDLE)
    add_text(s, "图：四类机组降容特性、风速—功率曲线与源荷缺口随温度的演化。",
             0.35, 6.68, 6.4, 0.5, 10, GRAY)

    # ---------- 2.3 基准 DC-OPF / MIQP 模型 ----------
    s = new_slide(prs)
    header(s, "2.3", "基准运行状态的 DC-OPF / MIQP 模型", 3)
    bold_title(s)
    label(s, "目标函数：发电成本 + 分级切负荷惩罚", 0.4, 1.30, 7.0, 13, DEEP)
    add_image_contain(s, EQ / "obj.png", 0.5, 1.66, 6.6, 0.55)
    label(s, "直流潮流与节点功率平衡", 0.4, 2.36, 6.5, 13, DEEP)
    add_image_contain(s, EQ / "balance.png", 0.5, 2.70, 6.6, 0.55)
    label(s, "火电启停与极热出力约束", 0.4, 3.40, 6.5, 13, DEEP)
    add_image_contain(s, EQ / "ucon.png", 0.5, 3.74, 6.4, 0.5)
    callout(s, 0.45, 4.42, 6.55, 2.25, [
        [("模型要点：", True)],
        [("• 含火电二次成本 + 0/1 启停 → ", False), ("混合整数二次规划（MIQP）", True)],
        [("• 基准 OPF 为极端气象后的静态快照，", False), ("不施加爬坡约束", True)],
        [("• 启停变量避免孤岛内", False), ("最小出力过剩导致的不可行", True)],
        [("• MATLAB/Gurobi 主求解 + Python 连续 QP 枚举", False), ("交叉验证", True)],
    ], size=12, valign=MSO_ANCHOR.MIDDLE)
    label(s, "基准调度结果与网络约束", 7.25, 1.30, 5.6, 13, DEEP)
    add_image_contain(s, FIG / "preview_03_dispatch_v2.png", 7.15, 1.66, 5.9, 3.55)
    callout(s, 7.20, 5.35, 5.85, 1.32, [
        [("关键发现：", True)],
        [("燃气机组 ", False), ("G4 极热上界 668.5 MW", True),
         ("，受支路热稳约束仅出力 ", False), ("651.9 MW", True)],
        [("→ 最优切负荷 ", False), ("762.81 MW", True),
         (" = 746.21 MW 缺口 + 16.60 MW 网络闲置", False)],
    ], size=11.5, valign=MSO_ANCHOR.MIDDLE)

    # ---------- 2.4 元件热故障概率模型 ----------
    s = new_slide(prs)
    header(s, "2.4", "极热条件下的元件热故障概率模型", 4)
    bold_title(s)
    callout(s, 0.45, 1.30, 5.15, 0.72, [
        [("统一框架：", True), ("应力相关故障率 λ → 评估窗口故障概率", False)],
        [("", False)],
    ], size=12, valign=MSO_ANCHOR.MIDDLE)
    add_image_contain(s, EQ / "pf.png", 5.75, 1.34, 2.5, 0.55)
    add_text(s, "（Δt = 24 h 暴露窗口）", 8.35, 1.46, 4.4, 0.3, 11, GRAY)
    blocks = [
        ("节点变压器", "trafo", "热点温度 + Arrhenius 绝缘热老化加速", RED,
         "环境温度驱动；DC-OPF 无无功，不用视在负载比"),
        ("发电机", "gen", "温度应力 + 出力率的比例风险模型", BLUE,
         "燃气机组温度敏感、出力率高，故障概率最高"),
        ("输电线路", "line", "IEEE 738 导线温度 + 过载应力放大", DEEP,
         "满载断面导线温度超连续运行限值，概率显著上升"),
    ]
    y = 2.25
    for name, eq, desc, color, note in blocks:
        add_rect(s, 0.45, y, 6.55, 1.42, WHITE, BOXLINE, True, 1.3)
        add_rect(s, 0.45, y, 0.10, 1.42, color)
        add_rich(s, [[(name, False)]], 0.66, y + 0.08, 2.2, 0.3, size=13.5, base=color)
        add_text(s, desc, 0.66, y + 0.40, 6.2, 0.3, 11.5, BLACK)
        add_image_contain(s, EQ / f"{eq}.png", 0.7, y + 0.70, 6.0, 0.46)
        add_text(s, "▶ " + note, 0.66, y + 1.15, 6.2, 0.24, 9.5, GRAY)
        y += 1.55
    add_image_contain(s, FIG / "preview_05_fault_probability_v2.png", 7.15, 1.90, 5.95, 4.05)
    callout(s, 7.20, 6.02, 5.85, 0.92, [
        [("24 h 条件故障概率：", True), ("G10 ≈ 4.11%、G4 ≈ 4.04%", False)],
        [("最高风险线路 L27 ≈ 0.42%；节点变压器 ≈ 0.015%", False)],
    ], size=11.5, valign=MSO_ANCHOR.MIDDLE)

    # ---------- 2.5 蒙特卡洛场景与故障后 OPF ----------
    s = new_slide(prs)
    header(s, "2.5", "蒙特卡洛故障场景与故障后 OPF 遍历", 5)
    bold_title(s)
    label(s, "① 85 维故障场景抽样", 0.4, 1.30, 5.0, 13, DEEP)
    add_image_contain(s, EQ / "sample.png", 0.5, 1.62, 3.2, 0.62)
    add_rich(s, [[
        ("按元件故障概率独立抽样，生成 ", False), ("2000 个 85 维", True),
        (" 0/1 场景（10 机 + 46 线 + 29 变压器）。", False),
    ]], 3.85, 1.62, 3.0, 0.95, size=11.5)
    label(s, "② 故障映射与拓扑重构", 6.95, 1.30, 5.8, 13, DEEP)
    add_rich(s, [
        [("• 发电机故障：", True), ("P(G)=0，火电强制停机", False)],
        [("• 线路故障：", True), ("支路退出，不参与 B 矩阵装配", False)],
        [("• 变压器故障：", True), ("关联支路退出，刻画通道受损", False)],
    ], 6.95, 1.62, 5.9, 1.05, size=11.5, line_spacing=1.15)
    add_line(s, 0.4, 2.85, 12.55, BOXLINE, 1.0)
    label(s, "③ 故障后 OPF：以基准出力为爬坡中心", 0.4, 2.95, 6.5, 13, DEEP)
    add_image_contain(s, EQ / "ramp.png", 0.5, 3.30, 6.2, 0.5)
    callout(s, 0.45, 3.95, 6.35, 2.75, [
        [("求解与统计结果：", True)],
        [("• 保留火电启停 → ", False), ("全部 2000 个场景均获最优解", True)],
        [("• 总故障事件 ", False), ("739", True), ("；含故障场景 ", False), ("631", True),
         ("；发电机事件占 ", False), ("70.2%", True)],
        [("• ", False), ("388", True), (" 个场景出现至少一台火电关停", False)],
        [("• 抽样经验频率与理论概率总体吻合，单场景最多 3 个故障", False)],
    ], size=11.5, valign=MSO_ANCHOR.MIDDLE)
    add_image_contain(s, FIG / "preview_06_monte_carlo_v2.png", 6.95, 3.05, 6.15, 3.7)

    # ---------- 2.6 失电风险评估结果 ----------
    s = new_slide(prs)
    header(s, "2.6", "失电风险评估结果与结论", 6)
    bold_title(s)
    add_image_contain(s, FIG / "preview_07_shed_distribution_v2.png", 0.35, 1.30, 7.35, 2.72)
    add_image_contain(s, FIG / "preview_08_risk_drivers_v2.png", 0.35, 4.10, 7.35, 2.62)
    add_image_contain(s, EQ / "bridge.png", 7.85, 1.30, 5.15, 0.66)
    kpis = [("951.58", "均值 / MW", BLUE), ("1765.56", "P95 / MW", DEEP), ("3330.40", "最大值 / MW", RED)]
    x = 7.85
    for val, lab, color in kpis:
        add_rect(s, x, 2.10, 1.62, 0.92, LIGHTBLUE, BOXLINE, True, 1.0)
        add_text(s, val, x + 0.03, 2.20, 1.56, 0.4, 17 if len(val) < 7 else 15, color, True, PP_ALIGN.CENTER)
        add_text(s, lab, x + 0.03, 2.63, 1.56, 0.24, 9.5, GRAY, False, PP_ALIGN.CENTER)
        x += 1.76
    callout(s, 7.85, 3.20, 5.15, 1.55, [
        [("风险分布特征：", True)],
        [("• 约 ", False), ("75%", True), (" 场景仍停在 762.81 MW 基准切负荷", False)],
        [("• 风险集中于", False), ("右侧尾部", True),
         ("，最大值约为基准的 ", False), ("4.4 倍", True)],
        [("• 电源故障后果强于单纯网络故障；", False),
         ("电源 + 网络混合故障风险最高", True)],
    ], size=11.5, valign=MSO_ANCHOR.MIDDLE)
    callout(s, 7.85, 4.92, 5.15, 1.78, [
        [("主要结论：", True)],
        [("① 极热无风形成 ", False), ("≈746 MW", True),
         (" 缺口，网络约束下切负荷 ", False), ("≈763 MW", True)],
        [("② 差异化 VOLL 使切负荷全部落在三级负荷", False)],
        [("③ 热故障叠加后切负荷呈", False), ("显著尾部风险", True)],
        [("④ 数值链条闭合并经 Python 交叉验证", False)],
    ], size=11.5, valign=MSO_ANCHOR.MIDDLE)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"written {OUT}")


if __name__ == "__main__":
    build()
